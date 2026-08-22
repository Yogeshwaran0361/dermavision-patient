import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    template_img = r'C:\Users\yoges\.gemini\antigravity\brain\56d6a50b-d761-49ef-a2a7-b9ec32a6764a\.user_uploaded\media_1786673854184.jpg'
    output_path = r'c:\Users\yoges\Downloads\techno\FINAL_POC_PRESENTATION.pptx'

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625) # 16:9 ratio matching template image

    blank_layout = prs.slide_layouts[6]

    # Color Palette
    DARK_BLUE = RGBColor(15, 23, 42)    # Slate 900
    EMERALD_GREEN = RGBColor(5, 150, 105) # Emerald 600
    SKY_BLUE = RGBColor(14, 165, 233)   # Sky 500
    TEXT_DARK = RGBColor(30, 41, 59)     # Slate 800
    TEXT_MUTED = RGBColor(100, 116, 139)# Slate 500
    ROSE_RED = RGBColor(225, 29, 72)

    def add_slide_bg(slide):
        if os.path.exists(template_img):
            slide.shapes.add_picture(template_img, 0, 0, width=prs.slide_width, height=prs.slide_height)

    def add_header(slide, title_text, category_text="DERMAVISION AI — 12-SLIDE POC PRESENTATION"):
        header_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.22), Inches(8.2), Inches(0.85))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p0 = tf.paragraphs[0]
        p0.text = category_text.upper()
        p0.font.size = Pt(9)
        p0.font.bold = True
        p0.font.color.rgb = SKY_BLUE

        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(19)
        p1.font.bold = True
        p1.font.color.rgb = DARK_BLUE

    # =============================================================
    # SLIDE 1: TITLE SLIDE
    # =============================================================
    s1 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s1)

    title_box = s1.shapes.add_textbox(Inches(1.6), Inches(1.1), Inches(8.0), Inches(4.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "KONGU ENGINEERING COLLEGE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = EMERALD_GREEN

    p = tf1.add_paragraph()
    p.text = "Department of Computer Science & Engineering"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED

    p = tf1.add_paragraph()
    p.text = "\nAI-Powered Skin Disease Analysis and Screening Assistance System"
    p.font.size = Pt(23)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    p = tf1.add_paragraph()
    p.text = "AI-Assisted Cutaneous Image Analysis with Patient Reporting & Tele-Dermatology Consultation"
    p.font.size = Pt(12.5)
    p.font.color.rgb = SKY_BLUE

    p = tf1.add_paragraph()
    p.text = "\nProof of Concept (POC) Technical Evaluation Presentation (12-Slide Comprehensive Deck)"
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK

    p = tf1.add_paragraph()
    p.text = "Core Stack: PyTorch ResNet50 Deep AI | OpenCV Quality Checker | FastAPI Backend | Firebase | React"
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED

    # =============================================================
    # SLIDE 2: PROBLEM STATEMENT & OBJECTIVES
    # =============================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s2)
    add_header(s2, "Problem Statement & Project Objectives")

    # Left Box: Problem Statement
    box2_left = s2.shapes.add_textbox(Inches(1.5), Inches(1.15), Inches(4.0), Inches(4.1))
    tf2_l = box2_left.text_frame
    tf2_l.word_wrap = True

    p = tf2_l.paragraphs[0]
    p.text = "PROBLEM STATEMENT"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ROSE_RED
    p.space_after = Pt(4)

    problems = [
        ("Rural Accessibility Barriers:", " People living in rural and remote areas cannot access certified dermatologists or specialized skin clinics due to long travel distances, lack of specialists, and healthcare costs."),
        ("Cancerous vs. Non-Cancerous Confusion:", " Patients cannot visually identify or differentiate whether an affected skin spot is a dangerous cancerous lesion (e.g. Malignant Melanoma / Actinic Keratosis) or a harmless non-cancerous mark (e.g. Common Mole / Benign Tag)."),
        ("Delayed Medical Intervention:", " Due to inability to recognize cancerous signs early, rural and semi-urban patients delay critical oncology consultations until advanced, life-threatening stages."),
        ("Lack of Immediate Screening Triage:", " Existing systems lack a simple digital screening tool to give instant risk stratification (Low vs Critical Cancer Risk) and connect patients with doctors.")
    ]

    for title, desc in problems:
        p = tf2_l.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(9.5)
        r1.font.color.rgb = DARK_BLUE
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(4)

    # Right Box: Objectives
    box2_right = s2.shapes.add_textbox(Inches(5.7), Inches(1.15), Inches(4.0), Inches(4.1))
    tf2_r = box2_right.text_frame
    tf2_r.word_wrap = True

    p = tf2_r.paragraphs[0]
    p.text = "PROJECT OBJECTIVES"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = EMERALD_GREEN
    p.space_after = Pt(4)

    objectives = [
        ("1. PyTorch ResNet50 Classifier:", " Build an automated AI engine classifying 10 skin disease categories."),
        ("2. OpenCV Image Validation:", " Pre-screen images for blur and brightness before AI execution."),
        ("3. Risk Stratification:", " Provide 4 risk levels (Low, Moderate, High, Critical) with Softmax confidence."),
        ("4. Tri-Lingual Reporting:", " Generate reports in English, Tamil, and Hindi with PDF export."),
        ("5. Firebase Cloud Hosting:", " Securely store profiles, scan history, and photos via Firestore & Storage."),
        ("6. Tele-Health Doctor Portal:", " Provide doctor workspace for Rx issuing and Google Meet video calls.")
    ]
    for title, desc in objectives:
        p = tf2_r.add_paragraph()
        r1 = p.add_run()
        r1.text = title
        r1.font.bold = True
        r1.font.size = Pt(9.5)
        r1.font.color.rgb = DARK_BLUE
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(4)

    # =============================================================
    # SLIDE 3: METHODOLOGY & TECHNICAL PIPELINE
    # =============================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s3)
    add_header(s3, "Methodology & End-to-End System Pipeline")

    box3 = s3.shapes.add_textbox(Inches(1.5), Inches(1.1), Inches(8.2), Inches(4.2))
    tf3 = box3.text_frame
    tf3.word_wrap = True

    p = tf3.paragraphs[0]
    p.text = "Complete Technical Execution Pipeline:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE
    p.space_after = Pt(6)

    steps3 = [
        "1. DATASET PREPARATION: 10 Clinical Disease Classes (Acne, Actinic Keratosis, Melanoma, Nevus, Psoriasis, etc.)",
        "2. PREPROCESSING & NORMALIZATION: RGB Conversion, 224x224 Resizing, ImageNet Mean & Std Normalization",
        "3. MODEL TRAINING: PyTorch ResNet50 Backbone + Custom Classification Head + Cross-Entropy Loss + Adam",
        "4. BACKEND SERVING: FastAPI Server (Python 3.11) + OpenCV Quality Checker (Laplacian Blur & Exposure)",
        "5. PATIENT FRONTEND SCANNER: React 19 + TypeScript + Vite 6 + Photo Upload / Smartphone Camera Dropzone",
        "6. INFERENCE & REPORT ENGINE: Softmax Ensemble Probabilities -> Multilingual Medical Knowledge Lookup",
        "7. FIREBASE CLOUD MANAGEMENT: Auth (Google Sign-In), Firestore (NoSQL Docs), Storage (Image Download URLs)",
        "8. DOCTOR PORTAL & CONSULTATION: Doctor Dashboard + Prescriptions (Rx) + 3-Dots Case Delete + Google Meet"
    ]

    for step in steps3:
        p = tf3.add_paragraph()
        p.text = "  ➔  " + step
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(4)

    # =============================================================
    # SLIDE 4: DATASET & PYTORCH MODEL ARCHITECTURE
    # =============================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s4)
    add_header(s4, "Dataset & PyTorch ResNet50 Architecture")

    box4 = s4.shapes.add_textbox(Inches(1.5), Inches(1.15), Inches(8.2), Inches(4.1))
    tf4 = box4.text_frame
    tf4.word_wrap = True

    details4 = [
        ("Model Architecture:", " Deep Residual Neural Network (PyTorch torchvision.models.resnet50) with 50 layers."),
        ("Skip Connections:", " Solves vanishing gradient problem via shortcut paths y = F(x) + x for unattenuated gradient flow."),
        ("Output Classes (10):", " Acne/Rosacea, Actinic Keratosis, Benign Other, Eczema/Dermatitis, Melanoma, Nevus/Mole, Psoriasis, Seborrheic Keratosis, Tinea Fungal, Vascular Lesion."),
        ("Input Shape & Transforms:", " Tensor (1, 3, 224, 224) RGB | Normalized with Mean=[0.485, 0.456, 0.406], Std=[0.229, 0.224, 0.225]."),
        ("Classification Head:", " 2048-dim Average Pooling -> Dropout(0.3) -> Linear(2048 -> 10)."),
        ("Optimizer & Loss:", " Adam Optimizer (Learning Rate = 10^-4) with Categorical Cross-Entropy Loss L = -sum(y * log(p))."),
        ("Weight Checkpoints:", " real_skin_classifier_v1.pth, real_skin_classifier_v2.pth (~90 MB trained weights).")
    ]

    for title, desc in details4:
        p = tf4.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = DARK_BLUE

        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(4)

    # =============================================================
    # SLIDE 5: RUNTIME AI INFERENCE & OPENCV QUALITY CHECK
    # =============================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s5)
    add_header(s5, "Runtime AI Inference & OpenCV Quality Check")

    box5 = s5.shapes.add_textbox(Inches(1.5), Inches(1.15), Inches(8.2), Inches(4.1))
    tf5 = box5.text_frame
    tf5.word_wrap = True

    steps5 = [
        ("1. Image Upload & Request:", " Patient uploads skin photo -> React sends multipart/form-data request to POST /api/ai/predict."),
        ("2. OpenCV Quality Checker:", " quality_checker.py computes Laplacian variance var(grad^2 I) for focus (> 20.0) and mean brightness (10–245 range). Warns user if blurry or dark."),
        ("3. Preprocessing Pipeline:", " Image bytes converted to PIL RGB -> Resized to 224x224 -> Transformed to PyTorch tensor with ImageNet normalization."),
        ("4. ResNet50 Inference Pass:", " PyTorch models run forward pass under torch.no_grad() -> Computes class probability distribution."),
        ("5. Softmax Probability Ensemble:", " Softmax converts logits into 10 class probabilities P_i = exp(z_i) / sum(exp(z_j)) -> Ensemble averages scores across checkpoints."),
        ("6. Disease & Risk Stratification:", " Argmax selects top class -> Retrieves localized clinical description, risk tier, symptoms & doctor advice.")
    ]

    for title, desc in steps5:
        p = tf5.add_paragraph()
        r1 = p.add_run()
        r1.text = title
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = DARK_BLUE

        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(4)

    # =============================================================
    # SLIDE 6: SYSTEM ARCHITECTURE & TECHNOLOGY STACK
    # =============================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s6)
    add_header(s6, "System Technology Stack & Architecture")

    box6 = s6.shapes.add_textbox(Inches(1.5), Inches(1.15), Inches(8.2), Inches(4.1))
    tf6 = box6.text_frame
    tf6.word_wrap = True

    arch_layers = [
        ("Patient Mobile Web App:", " React 19 + TypeScript + Vite 6 + Tailwind CSS 4 (Responsive SPA with dark mode)"),
        ("Doctor Web Portal:", " React 18 + TypeScript + Vite 5 + Tailwind CSS 3 (Specialized Tele-Health Workspace)"),
        ("AI Backend API Server:", " Python 3.11 + FastAPI + Uvicorn ASGI Server + PyTorch + Torchvision + OpenCV"),
        ("Firebase Authentication:", " Google Sign-In & Email/Password with Firebase Auth UID security isolation"),
        ("Cloud Firestore Database:", " Real-time NoSQL collections (users, scans, consultations, messages, notifications)"),
        ("Firebase Cloud Storage:", " Secure image hosting for scanned lesion photos with tokenized HTTPS URLs")
    ]

    for title, desc in arch_layers:
        p = tf6.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = SKY_BLUE

        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(5)

    # =============================================================
    # SLIDE 7: PATIENT APP WORKFLOW & USER JOURNEY
    # =============================================================
    s7 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s7)
    add_header(s7, "Patient App Workflow & User Journey")

    box7 = s7.shapes.add_textbox(Inches(1.5), Inches(1.15), Inches(8.2), Inches(4.1))
    tf7 = box7.text_frame
    tf7.word_wrap = True

    p_flow = [
        ("Step 1: Authentication:", " Patient signs in via Google Sign-In or Email/Password -> AuthContext initializes unique Firebase UID."),
        ("Step 2: Skin Scanner Page:", " Patient opens Scanner (/scanner) -> Uploads image or captures photo using smartphone camera."),
        ("Step 3: OpenCV Pre-Check & AI Scan:", " OpenCV validates quality -> FastAPI executes PyTorch ResNet50 inference in < 250ms."),
        ("Step 4: Immediate Result & Risk Badge:", " Displays top class prediction, confidence score (94.2%), and risk badge (Low, Moderate, High, Critical)."),
        ("Step 5: Diagnostic Report & History:", " Generates tri-lingual report -> Saves record to Firestore (users/{uid}/scans) for historical access."),
        ("Step 6: Doctor Consultation Booking:", " Patient clicks 'Consult Doctor' -> Creates consultation record for specialist review.")
    ]

    for title, desc in p_flow:
        p = tf7.add_paragraph()
        r1 = p.add_run()
        r1.text = title
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = DARK_BLUE

        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(5)

    # =============================================================
    # SLIDE 8: TRI-LINGUAL DIAGNOSTIC REPORT SYSTEM
    # =============================================================
    s8 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s8)
    add_header(s8, "Tri-Lingual Diagnostic Report System")

    box8 = s8.shapes.add_textbox(Inches(1.5), Inches(1.15), Inches(8.2), Inches(4.1))
    tf8 = box8.text_frame
    tf8.word_wrap = True

    rep_features = [
        ("Disease-Specific Content:", " Dynamically loads clinical overview, key symptoms, risk causes, precautions, and doctor referral guidance based on top AI prediction."),
        ("Tri-Lingual System:", " Full multi-language localization supporting English (en), Tamil (ta), and Hindi (hi)."),
        ("Reactive Zero-Reload Translation:", " Switching language from header selector (🌐) translates active AI reports and UI text instantly without page refresh."),
        ("Risk Color Stratification:", " Color-coded risk indicators (Critical Risk = Rose, High Risk = Amber, Moderate = Cyan, Low Risk = Emerald)."),
        ("PDF Report Export:", " Users can download print-formatted PDF diagnostic reports styled in their selected language.")
    ]

    for title, desc in rep_features:
        p = tf8.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = EMERALD_GREEN

        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(5)

    # =============================================================
    # SLIDE 9: FIREBASE CLOUD ARCHITECTURE & DATA MANAGEMENT
    # =============================================================
    s9 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s9)
    add_header(s9, "Firebase Cloud Architecture & Data Security")

    box9 = s9.shapes.add_textbox(Inches(1.5), Inches(1.15), Inches(8.2), Inches(4.1))
    tf9 = box9.text_frame
    tf9.word_wrap = True

    fb_items = [
        ("Firebase Authentication:", " Manages user credentials with OAuth 2.0 Google Sign-In & Email/Password. Enforces UID-based security rules."),
        ("Cloud Firestore Database:", " NoSQL hierarchical document collection architecture:"),
        ("   - users/{userId}:", " Patient profile data & language preferences."),
        ("   - users/{userId}/scans/{scanId}:", " Saved AI scan reports & historical predictions."),
        ("   - consultations/{id}:", " Consultation records, doctor diagnosis, and digital prescriptions (Rx:)."),
        ("   - notifications/{id}:", " Real-time patient invitation banners (restricted by patientUid rule)."),
        ("Firebase Cloud Storage:", " Hosts uploaded skin lesion images securely with tokenized HTTPS download URLs.")
    ]

    for title, desc in fb_items:
        p = tf9.add_paragraph()
        r1 = p.add_run()
        r1.text = title
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = DARK_BLUE

        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(4)

    # =============================================================
    # SLIDE 10: DOCTOR PORTAL & GOOGLE MEET TELE-HEALTH WORKFLOW
    # =============================================================
    s10 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s10)
    add_header(s10, "Doctor Portal & Google Meet Tele-Health Workflow")

    box10 = s10.shapes.add_textbox(Inches(1.5), Inches(1.15), Inches(8.2), Inches(4.1))
    tf10 = box10.text_frame
    tf10.word_wrap = True

    doc_items = [
        ("Separated Task Queues:", " Doctor Dashboard separates consultation queue into Pending Tasks (awaiting review) and Completed Tasks."),
        ("Structured Clinical Report Card:", " Displays patient scan specimen, AI confidence score, symptoms note, and clinical telemetry."),
        ("Digital Prescription Issuing:", " Doctors write official clinical diagnosis and prescription instructions (Rx:), saved to Firestore."),
        ("Google Meet Video Call:", " Doctor clicks 'Launch Google Meet' -> Starts call and sends real-time meeting link to target patient."),
        ("Targeted Patient Notifications:", " Real-time Firestore snapshot listener triggers a floating banner ('🎥 Doctor Launched Google Meet Call!') on patient screen."),
        ("Three Dots Case Deletion:", " Doctors can clean up old consultation cases via 3-dots case menu ('Delete Case').")
    ]

    for title, desc in doc_items:
        p = tf10.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = SKY_BLUE

        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(5)

    # =============================================================
    # SLIDE 11: EXPECTED OUTCOMES, INNOVATION & RESULTS
    # =============================================================
    s11 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s11)
    add_header(s11, "Expected Outcomes & Working Implementation Results")

    box11 = s11.shapes.add_textbox(Inches(1.5), Inches(1.15), Inches(8.2), Inches(4.1))
    tf11 = box11.text_frame
    tf11.word_wrap = True

    outcomes11 = [
        ("✔ Accelerated Triage & Screening:", " Sub-second preliminary skin lesion classification reduces patient panic and encourages early care."),
        ("✔ Risk Stratification Awareness:", " Flags high-risk malignant lesions (Melanoma, Actinic Keratosis) requiring urgent clinical attention."),
        ("✔ Enhanced Multi-Lingual Access:", " Tri-lingual support (English, Tamil, Hindi) ensures accessibility across regional backgrounds."),
        ("✔ Implemented Features Verified:", " PyTorch ResNet50 engine, OpenCV quality checker, React Patient SPA, Doctor Workspace, Firebase Firestore/Storage fully functional."),
        ("⚠️ Non-Diagnostic Disclaimer:", " System is designed as an AI screening assistance decision-support tool and does not replace certified dermatologists.")
    ]

    for title, desc in outcomes11:
        p = tf11.add_paragraph()
        r1 = p.add_run()
        r1.text = title
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = EMERALD_GREEN if "✔" in title else ROSE_RED

        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(5)

    # =============================================================
    # SLIDE 12: LIMITATIONS, FUTURE SCOPE & CONCLUSION
    # =============================================================
    s12 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s12)
    add_header(s12, "Limitations, Future Scope & Conclusion")

    box12 = s12.shapes.add_textbox(Inches(1.5), Inches(1.15), Inches(8.2), Inches(4.1))
    tf12 = box12.text_frame
    tf12.word_wrap = True

    p = tf12.paragraphs[0]
    p.text = "Limitations:"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ROSE_RED
    p.space_after = Pt(2)

    p = tf12.add_paragraph()
    p.text = "• Model scope bound to 10 clinical target classes; CPU inference latency ~100-200ms; screening assistance bounds."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(4)

    p = tf12.add_paragraph()
    p.text = "Future Scope & Enhancements:"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = EMERALD_GREEN
    p.space_after = Pt(2)

    p = tf12.add_paragraph()
    p.text = "➔ Grad-CAM Explainable AI visual heatmaps, WebRTC native video calls, expanded Fitzpatrick skin tone datasets."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(6)

    p = tf12.add_paragraph()
    p.text = "Conclusion:"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.space_after = Pt(2)

    p = tf12.add_paragraph()
    p.text = "DermaVision AI unites PyTorch Deep Learning + OpenCV Quality Defense + Tri-Lingual Reporting + Firebase Infrastructure + Doctor Tele-Health, making preliminary skin screening accessible, intelligent, and safely connected to medical professionals."
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE
    p.space_after = Pt(8)

    p_final = tf12.add_paragraph()
    p_final.text = "Thank You! Questions & Discussion."
    p_final.font.size = Pt(16)
    p_final.font.bold = True
    p_final.alignment = PP_ALIGN.CENTER
    p_final.font.color.rgb = SKY_BLUE

    try:
        prs.save(output_path)
        print("[SUCCESS] Successfully created 12-Slide PowerPoint Presentation at:", output_path)
    except Exception as e:
        fallback_path = r'c:\Users\yoges\Downloads\techno\FINAL_POC_PRESENTATION_12_SLIDES.pptx'
        prs.save(fallback_path)
        print(f"[SUCCESS NOTICE] Main file open in PowerPoint. Saved 12-Slide presentation to: {fallback_path}")

if __name__ == '__main__':
    create_presentation()

